"""
Script to create PowerPoint presentations for teaching demos.
Topics: Word Embeddings (5-min), PageRank (30-min), HashMap (30-min)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Alias for convenience
RgbColor = RGBColor

# Ensure output directory exists
os.makedirs('slides', exist_ok=True)

def set_slide_background(slide, r, g, b):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RgbColor(r, g, b)

def add_title_slide(prs, title, subtitle, color=(41, 128, 185)):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, *color)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(236, 240, 241)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, note=""):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RgbColor(44, 62, 80)
    
    # Add line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = RgbColor(41, 128, 185)
    line.line.fill.background()
    
    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(24)
        p.font.color.rgb = RgbColor(52, 73, 94)
        p.space_before = Pt(12)
    
    # Speaker notes
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note
    
    return slide

def add_quote_slide(prs, quote, author=""):
    """Add a quote slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 44, 62, 80)
    
    # Quote
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2.5))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(32)
    p.font.italic = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    if author:
        auth_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
        tf = auth_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"— {author}"
        p.font.size = Pt(24)
        p.font.color.rgb = RgbColor(189, 195, 199)
        p.alignment = PP_ALIGN.RIGHT
    
    return slide

def add_code_slide(prs, title, code):
    """Add a slide with code."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(44, 62, 80)
    
    # Code box background
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RgbColor(39, 40, 34)  # Dark background
    code_bg.line.fill.background()
    
    # Code text
    code_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(4.6))
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(16)
    p.font.name = "Consolas"
    p.font.color.rgb = RgbColor(248, 248, 242)  # Light text
    
    return slide

def add_equation_slide(prs, title, equation, explanation=""):
    """Add a slide with an equation."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RgbColor(44, 62, 80)
    
    # Equation (centered, large)
    eq_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = eq_box.text_frame
    p = tf.paragraphs[0]
    p.text = equation
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(41, 128, 185)
    p.alignment = PP_ALIGN.CENTER
    
    if explanation:
        exp_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
        tf = exp_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = explanation
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(52, 73, 94)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(44, 62, 80)
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RgbColor(39, 174, 96)
    
    # Left items
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(52, 73, 94)
        p.space_before = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.2), Inches(4.2), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RgbColor(231, 76, 60)
    
    # Right items
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4.2), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(52, 73, 94)
        p.space_before = Pt(8)
    
    return slide


# ============================================================
# WORD EMBEDDINGS - 5 Minute Video
# ============================================================

def create_word_embeddings_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 
                   "Word Embeddings", 
                   "Teaching Machines the Meaning of Words",
                   color=(155, 89, 182))
    
    # Slide 2: Images vs Language
    add_comparison_slide(prs, "Representing Data in Computers",
        "Images ✓", [
            "Naturally represented as matrices",
            "Pixels = numbers (0-255)",
            "28×28 image = 784 numbers",
            "Operations: convolution, pooling",
            "Easy for computers to process!"
        ],
        "Language ❓", [
            "Words are... symbols? Strings?",
            "How to convert 'cat' to numbers?",
            "How to capture meaning?",
            "How to show 'cat' ≈ 'dog'?",
            "This is the challenge!"
        ])
    
    # Slide 3: First Attempt - One-Hot Encoding
    add_content_slide(prs, "First Attempt: One-Hot Encoding", [
        "Vocabulary = {cat, dog, god, democracy, ...}  (10,000 words)",
        "",
        "'cat'  =  [0, 0, 0, 1, 0, 0, ..., 0]  ← position 4",
        "'dog'  =  [0, 1, 0, 0, 0, 0, ..., 0]  ← position 2",
        "'god'  =  [0, 0, 1, 0, 0, 0, ..., 0]  ← position 3",
        "'democracy'  =  [1, 0, 0, 0, 0, 0, ..., 0]  ← position 1",
        "",
        "Each word = 10,000-dimensional vector (sparse, one 1, rest 0s)"
    ])
    
    # Slide 4: The Distance Problem
    add_content_slide(prs, "The Distance Problem", [
        "Euclidean distance between any two one-hot vectors:",
        "",
        "Distance(cat, dog) = √2 ≈ 1.414",
        "Distance(cat, god) = √2 ≈ 1.414",
        "Distance(cat, democracy) = √2 ≈ 1.414",
        "",
        "Every word is EQUALLY different from every other! ❌",
        "No notion of similarity or meaning!"
    ])
    
    # Slide 5: Character-Level Similarity Problem
    add_content_slide(prs, "What About Character Similarity?", [
        "'dog' vs 'cat':  Semantically similar (both animals, pets)",
        "                 Character similarity: 0/3 letters match",
        "",
        "'dog' vs 'god':  Semantically different (pet vs deity)",
        "                 Character similarity: 2/3 letters match!",
        "",
        "Character-level similarity ≠ Semantic similarity",
        "We need meaning-based representations!"
    ])
    
    # Slide 6: What About N-grams?
    add_content_slide(prs, "What About N-grams?", [
        "Try bigrams: 'dog' = ['do', 'og']  'cat' = ['ca', 'at']",
        "             'god' = ['go', 'od']",
        "",
        "Still doesn't help with meaning:",
        "• 'dog' and 'god' share no bigrams, but neither do 'dog' and 'cat'",
        "• 'bank' (river) vs 'bank' (money) → Same n-grams! ❌",
        "",
        "N-grams capture orthography, not semantics"
    ])
    
    # Slide 7: The Key Insight - Quote
    add_quote_slide(prs, 
                   "You shall know a word by the company it keeps.",
                   "J.R. Firth, 1957")
    
    # Slide 8: Context Creates Meaning
    add_content_slide(prs, "The Distributional Hypothesis", [
        "'The dog barked loudly'      'The cat meowed softly'",
        "'Feed the dog treats'        'Feed the cat food'",
        "'Walk the dog daily'         'Pet the cat gently'",
        "",
        "dog and cat appear in similar contexts",
        "→ They should have similar meanings!",
        "→ They should have similar representations!"
    ])
    
    # Slide 9: Word Embeddings Solution
    add_content_slide(prs, "Word Embeddings: Dense Vectors", [
        "Learn dense vectors (e.g., 300 dimensions) from context",
        "",
        "'cat'  →  [0.2, -0.4, 0.7, ..., 0.5]  (300 numbers)",
        "'dog'  →  [0.3, -0.3, 0.6, ..., 0.4]  ← Close to cat!",
        "'god'  →  [-0.1, 0.8, -0.2, ..., -0.6]",
        "'democracy'  →  [-0.8, 0.1, -0.2, ..., 0.3]  ← Far from dog/cat",
        "",
        "Cosine similarity(dog, cat) ≈ 0.85   (very similar!)",
        "Cosine similarity(dog, god) ≈ 0.12   (not similar)"
    ])
    
    # Slide 10: How Are They Learned?
    add_content_slide(prs, "How Do We Learn Embeddings?", [
        "Word2Vec (2013): Predict context from word (or vice versa)",
        "",
        "Given: 'The quick brown fox jumps'",
        "Task: Predict 'brown' and 'jumps' from 'fox'",
        "",
        "Train on billions of sentences →",
        "Words in similar contexts get similar vectors!",
        "",
        "No manual labeling needed - fully unsupervised!"
    ])
    
    # Slide 11: Vector Arithmetic Magic
    add_equation_slide(prs, 
                      "The Magic of Vector Arithmetic",
                      "King - Man + Woman ≈ Queen",
                      "Relationships are encoded as directions in the vector space!\n\nParis - France + Italy ≈ Rome\nWalking - Walk + Swim ≈ Swimming")
    
    # Slide 12: Why It Works
    add_content_slide(prs, "Why Does Vector Arithmetic Work?", [
        "Gender direction: King - Man = Queen - Woman",
        "             → The 'male→female' vector is consistent!",
        "",
        "Capital direction: Paris - France = Rome - Italy",
        "              → The 'country→capital' vector is consistent!",
        "",
        "These relationships emerge from patterns in language",
        "Embeddings capture not just meaning, but relationships!"
    ])
    
    # Slide 13: Applications
    add_content_slide(prs, "Why It Matters", [
        "✓ Semantic similarity: 'happy' ≈ 'joyful' ≈ 'cheerful'",
        "✓ Better search: Query 'laptop' finds 'computer', 'notebook'",
        "✓ Machine translation: Align word spaces across languages",
        "✓ Sentiment analysis: Learn 'great' is positive",
        "✓ Recommendation systems: Similar items in vector space",
        "",
        "Foundation for modern NLP: BERT, GPT, all LLMs"
    ])
    
    # Slide 14: Evolution
    add_content_slide(prs, "The Evolution", [
        "2013: Word2Vec → Static embeddings",
        "      Same word, always same vector",
        "",
        "2018: ELMo, BERT → Contextual embeddings",
        "      'bank' (river) vs 'bank' (money) → Different vectors!",
        "",
        "2020+: GPT-3, ChatGPT → Even richer representations",
        "       But the core idea remains: meaning from context"
    ])
    
    # Slide 15: Key Takeaway
    add_title_slide(prs,
                   "Key Takeaway",
                   "From symbols to geometry:\nWords become points in space where meaning = proximity",
                   color=(46, 204, 113))
    
    # Slide 16: Thank You
    add_title_slide(prs,
                   "Thank You!",
                   "Questions?",
                   color=(155, 89, 182))
    
    prs.save('slides/Word_Embeddings_5min.pptx')
    print("✅ Created: Word_Embeddings_5min.pptx")


# ============================================================
# PAGERANK - 30 Minute Class
# ============================================================

def create_pagerank_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs,
                   "PageRank",
                   "How Google Organized the World's Information",
                   color=(41, 128, 185))
    
    # Slide 2: The Pre-Google Problem
    add_content_slide(prs, "The Year is 1997...", [
        "Internet has ~100 million web pages",
        "You search for 'jaguar' on AltaVista",
        "",
        "Results: 🚗 Car websites, 🐆 Animal pages, 💻 Mac OS",
        "         🏈 Sports teams, 📧 Spam, spam, spam...",
        "",
        "Keyword matching = Easily gamed, unreliable"
    ])
    
    # Slide 3: The Question
    add_content_slide(prs, "The Fundamental Question", [
        "",
        "How do you measure which page is",
        "genuinely IMPORTANT?",
        "",
        "— Larry Page & Sergey Brin, Stanford 1998"
    ])
    
    # Slide 4: Citation Analogy
    add_content_slide(prs, "The Academic Insight", [
        "How do we judge importance of research papers?",
        "→ Count citations!",
        "",
        "But not all citations are equal...",
        "A citation from Nature > citation from unknown journal",
        "",
        "Key Twist: Citation FROM important paper counts MORE"
    ])
    
    # Slide 5: Web as Graph
    add_content_slide(prs, "The Web as a Graph", [
        "📄 Every web page = Node",
        "🔗 Every hyperlink = Directed edge (a 'vote')",
        "",
        "A page is important if it's linked by important pages",
        "",
        "Notice: This is recursive! 🔄"
    ])
    
    # Slide 6: Fun Fact
    add_content_slide(prs, "Fun Fact", [
        "",
        "'PageRank' is a pun!",
        "",
        "📄  Ranks web PAGES",
        "👤  Named after Larry PAGE",
        ""
    ])
    
    # Slide 7: Random Surfer
    add_content_slide(prs, "The Random Surfer Model", [
        "Imagine a random web surfer (Alice)...",
        "",
        "1️⃣  Start on a random page",
        "2️⃣  Randomly click one of the links on the page",
        "3️⃣  Repeat... for a VERY long time",
        "",
        "Question: After billions of clicks,",
        "where is Alice most likely to be found?"
    ])
    
    # Slide 8: Random Surfer Answer
    add_content_slide(prs, "PageRank = Visit Probability", [
        "The answer gives us PageRank!",
        "",
        "Pages where Alice spends more time",
        "= More 'central' in the link structure",
        "= Higher PageRank",
        "",
        "[VISUAL: Animation of random walk converging]"
    ])
    
    # Slide 9: Damping Factor
    add_content_slide(prs, "The Damping Factor Problem", [
        "Problem 1: What if Alice hits a dead end? (no outlinks)",
        "Problem 2: What if there's a cycle she can't escape?",
        "",
        "Solution: With probability 15%, Alice gets 'bored'",
        "          and teleports to a random page!",
        "",
        "d = 0.85 (damping factor = probability of following a link)"
    ])
    
    # Slide 10: The Equation
    add_equation_slide(prs,
                      "The PageRank Equation",
                      "PR(A) = (1-d)/N + d × Σ PR(Ti)/L(Ti)",
                      "N = total pages    d = 0.85 (damping factor)\nM(A) = pages linking to A    L(Ti) = outlinks from Ti")
    
    # Slide 11: Intuition
    add_content_slide(prs, "Understanding the Formula", [
        "(1-d)/N  →  Teleportation (random jump probability)",
        "",
        "d × Σ PR(Ti)/L(Ti)  →  Sum of 'votes' from linking pages",
        "",
        "If a page has PR=0.5 and 10 outlinks,",
        "each link passes 0.05 of its importance"
    ])
    
    # Slide 12: Example
    add_content_slide(prs, "Worked Example: 4 Pages", [
        "A ← linked by B, C, D",
        "B ← linked by D",
        "C ← linked by B, D",
        "",
        "Start: All pages have PR = 0.25",
        "Iterate until convergence...",
        "",
        "Result: A has highest PageRank (most incoming links)"
    ])
    
    # Slide 13: Matrix Form
    add_equation_slide(prs,
                      "Matrix Formulation",
                      "r = M × r",
                      "PageRank vector r is the principal eigenvector\nof the transition matrix M with eigenvalue 1\n\nSolution: Power iteration - multiply until convergence")
    
    # Slide 14: Convergence
    add_content_slide(prs, "Convergence is Fast!", [
        "Original Google (1998):",
        "• 322 million pages",
        "• Converged in just 52 iterations",
        "",
        "The algorithm scales beautifully ✓"
    ])
    
    # Slide 15: Code
    add_code_slide(prs, "PageRank Algorithm (Python)", 
"""def pagerank(M, d=0.85, epsilon=1e-8):
    N = M.shape[0]
    r = np.ones(N) / N  # Start uniform
    
    while True:
        r_new = (1-d)/N + d * M @ r
        if np.linalg.norm(r_new - r) < epsilon:
            return r_new
        r = r_new""")
    
    # Slide 16: Edge Cases
    add_content_slide(prs, "Handling Edge Cases", [
        "🔹 Dangling nodes (no outlinks):",
        "   Distribute their PageRank evenly to all pages",
        "",
        "🔹 Spider traps (self-loops, cycles):",
        "   Damping factor prevents getting stuck",
        "",
        "🔹 Sparse matrix storage:",
        "   Most pages link to only a few others"
    ])
    
    # Slide 17: Applications Beyond Web
    add_content_slide(prs, "Beyond Web Search", [
        "PageRank works on ANY directed graph!",
        "",
        "📚 Academic citations → Measure paper influence",
        "🐦 Twitter 'Who to Follow' → Personalized PageRank",
        "🧬 Biology → Essential proteins in networks",
        "🏈 Sports → Team rankings from game results"
    ])
    
    # Slide 18: Limitations
    add_content_slide(prs, "Limitations", [
        "❌ Topic-agnostic: High-PR cat page won't help dog search",
        "❌ Favors old pages: New pages haven't accumulated links",
        "❌ Gaming: Link farms tried to manipulate rankings",
        "",
        "Today: Google uses 200+ ranking signals",
        "But PageRank remains a foundational signal"
    ])
    
    # Slide 19: Key Takeaways
    add_content_slide(prs, "Key Takeaways", [
        "1️⃣ Important pages are linked by important pages",
        "2️⃣ Random surfer: PageRank = visit probability",
        "3️⃣ Math: Principal eigenvector of transition matrix",
        "4️⃣ Fast convergence, scales to billions of pages",
        "5️⃣ Idea extends far beyond web search"
    ])
    
    # Slide 20: Bigger Picture
    add_content_slide(prs, "The Bigger Picture", [
        "",
        "From content-based → structure-based analysis",
        "",
        "The links between things often tell us",
        "more than the things themselves",
        ""
    ])
    
    # Slide 21: Questions
    add_title_slide(prs,
                   "Questions?",
                   "PageRank: Turning chaos into organized information",
                   color=(41, 128, 185))
    
    prs.save('slides/PageRank_30min.pptx')
    print("✅ Created: PageRank_30min.pptx")


# ============================================================
# HASH TABLES / HASHMAP - 30 Minute Class
# ============================================================

def create_hashmap_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs,
                   "Hash Tables",
                   "The Secret Behind Lightning-Fast Lookups",
                   color=(230, 126, 34))
    
    # Slide 2: The Challenge
    add_content_slide(prs, "The Challenge", [
        "📞 Phonebook with 1 million entries",
        "Find: 'Sharma, Priya' → Phone number",
        "",
        "Unsorted list: Check all 1 million → O(n)",
        "Sorted list + Binary search: ~20 comparisons → O(log n)",
        "",
        "Python dict: ONE step → O(1) 🚀"
    ])
    
    # Slide 3: Where Hash Tables Live
    add_content_slide(prs, "Hash Tables Power Everything", [
        "🐍 Python: dict, set",
        "☕ Java: HashMap, HashSet",
        "📜 JavaScript: Objects",
        "🗄️ Databases: Indexes",
        "💾 Caching systems",
        "🌐 Your browser's local storage"
    ])
    
    # Slide 4: Why Arrays are Fast
    add_content_slide(prs, "Arrays: The Fastest Data Structure", [
        "Array[42] → Direct address calculation:",
        "   address = base + 42 × element_size",
        "",
        "ONE step! O(1) ✓",
        "",
        "But... arrays need INTEGER indices",
        "What about string keys like names? 🤔"
    ])
    
    # Slide 5: Hash Function Idea
    add_content_slide(prs, "The Brilliant Idea: Hash Functions", [
        "",
        "What if we could convert ANY key → integer?",
        "",
        "'Priya'  →  hash()  →  7",
        "'Rahul'  →  hash()  →  42",
        "",
        "Then use that integer as an array index!"
    ])
    
    # Slide 6: Hash Table in Action
    add_content_slide(prs, "Hash Table Operations", [
        "INSERT 'Priya' → hash('Priya') = 7 → store at index 7",
        "",
        "LOOKUP 'Priya' → hash('Priya') = 7 → retrieve index 7",
        "",
        "Both operations: O(1)! 🎉"
    ])
    
    # Slide 7: Hash Function Example
    add_code_slide(prs, "Simple Hash Function",
"""def simple_hash(key, table_size):
    total = 0
    for char in key:
        total = (total * 31 + ord(char))
    return total % table_size

# Example:
hash('Priya', 10)  # → 7
hash('Rahul', 10)  # → 3""")
    
    # Slide 8: The Collision Problem
    add_content_slide(prs, "The Inevitable Problem: Collisions", [
        "Infinite possible keys (all strings)",
        "Finite array slots (e.g., 100)",
        "",
        "By pigeonhole principle:",
        "Multiple keys MUST map to the same index!",
        "",
        "hash('Priya') = 7",
        "hash('Alicia') = 7   ← COLLISION! 💥"
    ])
    
    # Slide 9: Chaining
    add_content_slide(prs, "Solution 1: Chaining", [
        "At each index, store a linked list",
        "",
        "Index 7: → [Priya: 555-1234] → [Alicia: 555-5678]",
        "",
        "INSERT: Hash → Append to list",
        "LOOKUP: Hash → Search through list",
        "DELETE: Hash → Remove from list"
    ])
    
    # Slide 10: Open Addressing
    add_content_slide(prs, "Solution 2: Open Addressing", [
        "No linked lists. If slot is taken, probe for next empty.",
        "",
        "Linear probing: try i, i+1, i+2, ...",
        "",
        "Insert 'Priya' → hash=7 → slot 7 empty → ✓",
        "Insert 'Alicia' → hash=7 → slot 7 full → try 8 → ✓",
        "Lookup 'Alicia' → hash=7 → not Alicia → try 8 → found!"
    ])
    
    # Slide 11: Probing Strategies
    add_content_slide(prs, "Open Addressing Variants", [
        "Linear probing: i, i+1, i+2, ...",
        "   Simple but causes clustering",
        "",
        "Quadratic probing: i, i+1², i+2², ...",
        "   Reduces clustering",
        "",
        "Double hashing: i, i+h₂(key), i+2h₂(key), ...",
        "   Best distribution"
    ])
    
    # Slide 12: Comparison
    add_comparison_slide(prs, "Chaining vs Open Addressing",
        "Chaining ✓", [
            "Simpler implementation",
            "Load factor can exceed 1",
            "Easy deletion",
            "Used by: Java HashMap"
        ],
        "Open Addressing ✓", [
            "Better cache performance",
            "More compact memory",
            "Tricky deletion (tombstones)",
            "Used by: Python dict"
        ])
    
    # Slide 13: Load Factor
    add_content_slide(prs, "Load Factor: α = n/m", [
        "n = items stored",
        "m = table size",
        "",
        "As α grows → more collisions → slower",
        "",
        "Chaining: Still okay at α = 2 or 3",
        "Open addressing: Crashes as α → 1!"
    ])
    
    # Slide 14: Dynamic Resizing
    add_content_slide(prs, "Dynamic Resizing", [
        "When load factor exceeds threshold (typically 0.75):",
        "",
        "1️⃣ Create new array (2× the size)",
        "2️⃣ REHASH all existing elements",
        "3️⃣ Insert into new table",
        "",
        "Why rehash? hash depends on table size!"
    ])
    
    # Slide 15: Amortized Analysis
    add_content_slide(prs, "But Rehashing is O(n)?!", [
        "Yes, but we only resize after n insertions",
        "",
        "Amortized analysis:",
        "Each insert 'pays' for a piece of future resize",
        "",
        "Average per insert: Still O(1)! ✓"
    ])
    
    # Slide 16: Good Hash Function
    add_content_slide(prs, "What Makes a Good Hash Function?", [
        "✓ Deterministic: Same input → same output",
        "✓ Fast: O(key length) is ideal",
        "✓ Uniform distribution: All buckets equally likely",
        "✓ Avalanche effect: Similar keys → different indices"
    ])
    
    # Slide 17: Bad vs Good Hash
    add_comparison_slide(prs, "Hash Function Quality",
        "Bad Hash ❌", [
            "def bad(s): return ord(s[0])",
            "All 'A' names → same bucket!",
            "Terrible distribution",
            "Degrades to O(n)"
        ],
        "Good Hash ✓", [
            "Polynomial rolling hash",
            "h = s[0]×31⁰ + s[1]×31¹ + ...",
            "Excellent distribution",
            "Maintains O(1)"
        ])
    
    # Slide 18: Applications
    add_content_slide(prs, "Real-World Applications", [
        "🔄 Caching: O(1) lookup for cached results",
        "🗄️ Database indexing: Fast equality queries",
        "🔍 Deduplication: Find duplicates in O(n)",
        "📊 Counting: Word frequency in O(n)",
        "🔗 Symbol tables: Compilers, interpreters"
    ])
    
    # Slide 19: Word Count Example
    add_code_slide(prs, "Example: Word Counting",
"""# Without hash table: O(n²)
# With hash table: O(n) !

counts = {}  # Hash table

for word in text.split():
    counts[word] = counts.get(word, 0) + 1

# 'the': 42
# 'quick': 3
# 'brown': 5""")
    
    # Slide 20: In Different Languages
    add_content_slide(prs, "Hash Tables in Programming Languages", [
        "🐍 Python: dict, set → Open addressing",
        "☕ Java: HashMap → Chaining + trees",
        "🔧 C++: unordered_map → Chaining",
        "📜 JavaScript: Objects → Highly optimized"
    ])
    
    # Slide 21: Complexity Summary
    add_content_slide(prs, "Time Complexity", [
        "           Average    Worst Case",
        "Insert:      O(1)        O(n)",
        "Lookup:      O(1)        O(n)",
        "Delete:      O(1)        O(n)",
        "",
        "Worst case: All elements hash to same bucket",
        "(Astronomically unlikely with good hash function)"
    ])
    
    # Slide 22: Common Pitfalls
    add_content_slide(prs, "⚠️ Common Pitfalls", [
        "🔹 Mutable keys: If key changes, lost forever!",
        "🔹 Bad hash functions: Degrades to O(n)",
        "🔹 Hash collision attacks: Crafted inputs for O(n²)",
        "🔹 Java: Must override hashCode with equals"
    ])
    
    # Slide 23: Key Takeaways
    add_content_slide(prs, "Key Takeaways", [
        "1️⃣ O(1) average for insert, lookup, delete",
        "2️⃣ Trade space for time",
        "3️⃣ Two strategies: Chaining vs Open Addressing",
        "4️⃣ Resizing maintains performance (amortized O(1))",
        "5️⃣ Hash function quality is critical"
    ])
    
    # Slide 24: The Big Idea
    add_content_slide(prs, "The Bigger Picture", [
        "",
        "Clever indexing beats searching",
        "",
        "Instead of looking through data,",
        "compute exactly where it should be",
        ""
    ])
    
    # Slide 25: Questions
    add_title_slide(prs,
                   "Questions?",
                   "Hash Tables: O(1) magic at your fingertips",
                   color=(230, 126, 34))
    
    prs.save('slides/HashMap_30min.pptx')
    print("✅ Created: HashMap_30min.pptx")


# ============================================================
# DYNAMIC PROGRAMMING - 5 Minute Video
# ============================================================

def create_dynamic_programming_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs,
                   "Dynamic Programming",
                   "The Art of Being Lazy... Smartly",
                   color=(142, 68, 173))
    
    # Slide 2: Opening Hook
    add_content_slide(prs, "The Mystery", [
        "Calculate Fibonacci(40):",
        "   Naive recursion: 30 seconds ⏰",
        "",
        "Calculate Fibonacci(100):",
        "   With Dynamic Programming: 0.001 seconds ⚡",
        "",
        "What changed?"
    ])
    
    # Slide 3: The Problem - Fibonacci Tree
    add_content_slide(prs, "Why Is Naive Recursion Slow?", [
        "Computing Fib(5) recursively:",
        "",
        "              F(5)",
        "           /        \\",
        "        F(4)         F(3)    ← F(3) computed twice!",
        "       /    \\        /   \\",
        "    F(3)   F(2)   F(2) F(1) ← F(2) computed 3 times!",
        "",
        "For Fib(40): Over 1 BILLION redundant calculations! 🔥"
    ])
    
    # Slide 4: The DP Insight
    add_title_slide(prs,
                   "Dynamic Programming Asks:",
                   "Why solve the same problem twice?",
                   color=(52, 152, 219))
    
    # Slide 5: Three Principles
    add_content_slide(prs, "The Three Principles of DP", [
        "1️⃣  REMEMBER",
        "    Store solutions to subproblems",
        "",
        "2️⃣  REUSE",
        "    Look up instead of recomputing",
        "",
        "3️⃣  BUILD",
        "    Solve small problems first, combine for bigger ones"
    ])
    
    # Slide 6: Fibonacci Bottom-Up
    add_content_slide(prs, "Fibonacci with DP: Build from Bottom", [
        "Instead of recursing down, build up from base cases:",
        "",
        "F(0) = 0",
        "F(1) = 1",
        "F(2) = F(0) + F(1) = 1",
        "F(3) = F(1) + F(2) = 2",
        "F(4) = F(2) + F(3) = 3",
        "F(5) = F(3) + F(4) = 5",
        "",
        "Time: O(n) instead of O(2ⁿ) ✓"
    ])
    
    # Slide 7: The Transformation
    add_comparison_slide(prs, "The Transformation",
        "Naive Recursion ❌", [
            "Time: O(2ⁿ)",
            "Fib(40): ~30 seconds",
            "Fib(100): Heat death of universe",
            "Redundant calculations"
        ],
        "Dynamic Programming ✓", [
            "Time: O(n)",
            "Fib(40): <1 millisecond",
            "Fib(100): Still <1 millisecond",
            "Each subproblem solved once"
        ])
    
    # Slide 8: Real-World Problem - Edit Distance
    add_content_slide(prs, "Real-World Application: Auto-Correct", [
        "",
        "How does your phone know you meant",
        "'kitten' when you typed 'sitting'?",
        "",
        "Answer: Edit Distance!",
        ""
    ])
    
    # Slide 9: Edit Distance Definition
    add_content_slide(prs, "Edit Distance Problem", [
        "Minimum operations to transform one string into another",
        "",
        "Allowed operations:",
        "• Insert a character",
        "• Delete a character",
        "• Substitute a character",
        "",
        "Example: 'kitten' → 'sitting'"
    ])
    
    # Slide 10: Edit Distance Example
    add_content_slide(prs, "Example: 'kitten' → 'sitting'", [
        "Step 1: Substitute k → s:    'sitten'",
        "Step 2: Substitute e → i:    'sittin'",
        "Step 3: Insert g at end:     'sitting'",
        "",
        "Total: 3 operations",
        "",
        "Question: Is this the minimum? How do we know?"
    ])
    
    # Slide 11: DP Table Introduction
    add_content_slide(prs, "The DP Approach: Build a Table", [
        "dp[i][j] = edit distance between:",
        "• First i characters of 'kitten'",
        "• First j characters of 'sitting'",
        "",
        "Base cases:",
        "• dp[0][j] = j  (insert j characters)",
        "• dp[i][0] = i  (delete i characters)"
    ])
    
    # Slide 12: The Recurrence
    add_content_slide(prs, "The DP Recurrence", [
        "For each cell dp[i][j]:",
        "",
        "IF characters match (kitten[i] == sitting[j]):",
        "    dp[i][j] = dp[i-1][j-1]  (no operation needed)",
        "",
        "ELSE, take minimum of:",
        "    • dp[i-1][j] + 1      (delete from first string)",
        "    • dp[i][j-1] + 1      (insert into first string)",
        "    • dp[i-1][j-1] + 1    (substitute)"
    ])
    
    # Slide 13: Filled Table
    add_code_slide(prs, "The Complete DP Table", 
"""      ""  s  i  t  t  i  n  g
""     0  1  2  3  4  5  6  7
k      1  1  2  3  4  5  6  7
i      2  2  1  2  3  4  5  6
t      3  3  2  1  2  3  4  5
t      4  4  3  2  1  2  3  4
e      5  5  4  3  2  2  3  4
n      6  6  5  4  3  3  2  3
       
Bottom-right = 3 ← Our answer!""")
    
    # Slide 14: How to Fill
    add_content_slide(prs, "Building the Table", [
        "Fill row by row, left to right",
        "",
        "Example: dp[2][2] (comparing 'ki' with 'si')",
        "• kitten[2]='i', sitting[2]='i' → Match!",
        "• dp[2][2] = dp[1][1] = 1 ✓",
        "",
        "Example: dp[1][1] (comparing 'k' with 's')",
        "• No match, so take min(dp[0][1], dp[1][0], dp[0][0]) + 1",
        "• min(1, 1, 0) + 1 = 1"
    ])
    
    # Slide 15: Complexity Analysis
    add_content_slide(prs, "Time & Space Complexity", [
        "Without DP:",
        "• Try all possible edit sequences → Exponential! 💥",
        "",
        "With DP:",
        "• Time: O(n × m) where n, m are string lengths",
        "• Space: O(n × m) for the table",
        "",
        "For 'kitten' (6) and 'sitting' (7): Only 42 cells!",
        "Tractable for real applications ✓"
    ])
    
    # Slide 16: Applications
    add_content_slide(prs, "Where Is Edit Distance Used?", [
        "📱 Spell checkers & autocorrect",
        "   'teh' → 'the' (distance = 2)",
        "",
        "🔍 Search suggestions",
        "   'jagguar' → 'jaguar' (distance = 1)",
        "",
        "🧬 DNA sequence alignment (bioinformatics)",
        "   Comparing genetic sequences",
        "",
        "📄 Plagiarism detection",
        "   Measuring document similarity"
    ])
    
    # Slide 17: More DP Applications
    add_content_slide(prs, "Dynamic Programming Powers...", [
        "✓ Longest Common Subsequence (diff tools)",
        "✓ Matrix Chain Multiplication (GPU optimization)",
        "✓ Knapsack Problem (resource allocation)",
        "✓ Sequence Alignment (NLP, bioinformatics)",
        "✓ Optimal Binary Search Trees (databases)",
        "✓ And many, many more..."
    ])
    
    # Slide 18: Key Takeaway
    add_title_slide(prs,
                   "Key Takeaway",
                   "Break problems into overlapping subproblems\nSolve once. Remember forever.\n\nIt's not magic—it's smart bookkeeping!",
                   color=(46, 204, 113))
    
    # Slide 19: Challenge
    add_content_slide(prs, "Challenge for You!", [
        "",
        "Calculate the edit distance from",
        "YOUR name to your friend's name",
        "",
        "Build the DP table!",
        ""
    ])
    
    # Slide 20: Thank You
    add_title_slide(prs,
                   "Thank You!",
                   "Questions?",
                   color=(142, 68, 173))
    
    prs.save('slides/Dynamic_Programming_5min.pptx')
    print("✅ Created: Dynamic_Programming_5min.pptx")


# ============================================================
# MAIN
# ============================================================

if __name__ == "__main__":
    print("Creating PowerPoint presentations...\n")
    create_word_embeddings_pptx()
    create_pagerank_pptx()
    create_hashmap_pptx()
    create_dynamic_programming_pptx()
    print("\n✅ All presentations created in the 'slides' folder!")
